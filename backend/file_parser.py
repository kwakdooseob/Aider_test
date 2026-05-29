import docx
from pdfminer.high_level import extract_text
from pptx import Presentation
import os

def get_file_type(file_path):
    """파일 확장자를 기반으로 파일 타입을 판별합니다."""
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    if ext == '.docx':
        return 'docx'
    elif ext == '.pdf':
        return 'pdf'
    elif ext == '.pptx':
        return 'pptx'
    else:
        return None

def parse_docx(file_path):
    try:
        doc = docx.Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text.strip())
        return '\n'.join(full_text).strip()
    except Exception as e:
        print(f"Error parsing docx file: {e}")
        return None

def parse_pdf(file_path):
    try:
        text = extract_text(file_path)
        if text:
            return ' '.join(text.split()).strip()
        return ''
    except Exception as e:
        print(f"Error parsing pdf file: {e}")
        return None

def parse_pptx(file_path):
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text.strip())
        return '\n'.join(full_text).strip()
    except Exception as e:
        print(f"Error parsing pptx file: {e}")
        return None
